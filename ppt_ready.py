import json
import os
import re
import logging
from datetime import datetime
from PIL import Image

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from deep_translator import GoogleTranslator
# 引入配置文件
import config

# 配置日志
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

class PPTGenerator:
    def __init__(self, data, template_path, images_dir, location_name, language):
        self.data = data
        self.template_path = template_path
        self.images_dir = images_dir
        self.location = location_name
        self.language = language
        self.prs = None

    def load_resources(self):
        """加载模板和JSON数据"""
        try:
            if not os.path.exists(self.template_path):
                raise FileNotFoundError(f"模板文件未找到: {self.template_path}")
            self.prs = Presentation(self.template_path)
            # logging.info(f"模板加载成功，共有 {len(self.prs.slide_layouts)} 个布局")

            logging.info("资源加载成功")
        except Exception as e:
            logging.error(f"资源加载失败: {e}")
            raise

    # --- 通用工具方法 ---

    def _set_text_style(self, run, font_name='华文细黑', size=12, bold=False, color=config.COLOR_BLACK):
        """统一设置文本样式"""
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    def _get_image_dimensions(self, image_path):
        """获取图片尺寸"""
        try:
            with Image.open(image_path) as img:
                return img.size
        except Exception as e:
            logging.warning(f"无法读取图片尺寸 {image_path}: {e}")
            return 0, 0

    def _calculate_fitted_size(self, img_width, img_height, max_width, max_height):
        """计算适合给定区域的图片大小，保持纵横比"""
        if img_height == 0 or max_height == 0: return 0, 0
        
        aspect = img_width / img_height
        max_aspect = max_width / max_height

        if aspect > max_aspect:
            fitted_width = max_width
            fitted_height = max_width / aspect
        else:
            fitted_height = max_height
            fitted_width = max_height * aspect
        return fitted_width, fitted_height

    def _remove_all_picture_placeholders(self, slide):
        """移除幻灯片中所有的图片占位符"""
        for shape in list(slide.shapes):
            if shape.is_placeholder and shape.placeholder_format.type == 18:
                try:
                    sp = shape._element
                    sp.getparent().remove(sp)
                except Exception:
                    pass
    def _remove_all_text_placeholders(self, slide):
        """移除幻灯片中所有的文本占位符"""
        for shape in list(slide.shapes):
            if shape.is_placeholder and shape.placeholder_format.type == 15:
                try:
                    sp = shape._element
                    sp.getparent().remove(sp)
                except Exception:
                    pass

    def translate_with_glossary(self,text):
        term_dict = {
            "Bloomberg": "彭博",
            "Reuters": "路透",
            "Goldman Sachs": "高盛",
        }
    
        # 先替换已知术语
        for eng, zh in term_dict.items():
            text = text.replace(eng, zh)
        
        # # 剩余部分尝试翻译（如果有英文）
        # if any(c.isalpha() for c in text):
        #     try:
        #         return GoogleTranslator(source='en', target='zh-CN').translate(text)
        #     except:
        #         return text
        return text

    def _get_standard_keys(self, title):
        """
        根据文章标题中的关键词，返回该资产类别的【中文标准名称】。
        """
        title_map = {
            # 关键词 : 中文标准名
            "US Equities": "美股",
            "HK/China Equities": "中港股市",
            "European Equities": "欧股", 
            "Japan Equities": "日股",
            "Fixed Income": "债券",
            "Gold": "黄金",
            "Crude Oil": "原油",
            "Fund flow": "资金流",
            "Top Picks - Bonds": "个债精选",   
            "Top Picks - Equities": "个股精选",
            "债市": "债券"
        }
        
        # 统一转为小写进行匹配（不区分大小写）
        title_lower = title.lower()
        
        # 遍历映射表，如果标题包含关键词，就返回对应的中文标准名
        for keyword, chinese_name in title_map.items():
            keyword_lower = keyword.lower()
            if keyword_lower in title_lower:
                return chinese_name
        
        # 如果没有匹配到，返回原始标题的前2个字
        return title[:2]

    def _find_matching_image(self, title, images_dir):
        """
        图片匹配：先找到对应的中文标准名，再取前2个字进行匹配
        """
        if not os.path.exists(images_dir):
            logging.warning(f"图片目录不存在: {images_dir}")
            return None
        
        # 1. 先通过_get_standard_keys获取中文标准名
        chinese_name = self._get_standard_keys(title)
        
        # 2. 取中文标准名的前2个字作为键
        if len(chinese_name) >= 2:
            key = chinese_name[:2]
        else:
            key = chinese_name
        
        # 如果没有有效的键，返回None
        if not key or len(key.strip()) == 0:
            return None
        
        # 打印调试信息
        print(f"    查找图片，原始标题: '{title}'")
        print(f"    中文标准名: '{chinese_name}'")
        print(f"    使用键: '{key}'")
        
        # 支持的图片格式
        image_extensions = ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp']
        
        # 遍历图片目录，只匹配以键开头的文件名
        for filename in os.listdir(images_dir):
            # 检查是否为图片文件
            if any(filename.lower().endswith(ext.lower()) for ext in image_extensions):
                # 文件名以键开头（精确匹配）
                if filename.startswith(key):
                    print(f"    匹配成功: 键 '{key}' -> 图片 '{filename}'")
                    return os.path.join(images_dir, filename)
        
        print(f"    警告: 未找到以 '{key}' 开头的图片")
        return None

    def _fill_image(self, slide, image_path):
        """在幻灯片中插入图片并调整填充方式"""
        if not os.path.exists(image_path): return None

        img_width, img_height = self._get_image_dimensions(image_path)
        if img_width == 0: return None
        
        # 查找图片占位符
        placeholder = None
        for shape in slide.shapes:
            if shape.placeholder_format.type == 18:
                placeholder = shape
                break

        if placeholder:
            w, h = self._calculate_fitted_size(img_width, img_height, placeholder.width, placeholder.height)
            left = placeholder.left + (placeholder.width - w) // 2
            top = placeholder.top + (placeholder.height - h) // 2
            
            try:
                picture = slide.shapes.add_picture(image_path, left, top, w, h)
                logging.info(f"✓ 插入图片: {os.path.basename(image_path)}")
                return picture
            except Exception as e:
                logging.error(f"插入图片失败: {e}")
        return None

    def _add_image_annotations(self, slide, image_path):
        """添加图表标题和资料来源 (使用 config 中的坐标配置)"""
        try:
            file_name = os.path.basename(image_path)
            parts = os.path.splitext(file_name)[0].split('_')

           
            if len(parts) >= 2:
                chart_title = parts[1]
                
                # 1. 翻译逻辑 (如果是英文模式)
                is_english_mode = (self.language == 'en') # 标记是否为英文模式
                
                if chart_title and chart_title != "NONE" and is_english_mode:
                    try:
                        chart_title = GoogleTranslator(source='auto', target='en').translate(chart_title)
                    except Exception as e:
                        print(f"Translation failed: {e}")

                if chart_title and chart_title != "NONE":
                    cfg = config.ANNOTATION_CONFIG['title']
                    
                    # 2. 动态调整左边距 (英文和中文长度判断标准不同)
                    offset = Pt(0)
                    print(f"length of chart title '{chart_title}': {len(chart_title)}")
                    if is_english_mode:
                        # 英文判断逻辑 (字符数较多，但 Arial Narrow 比较省空间)
                        if len(chart_title) > 75: offset = Pt(140)
                        elif len(chart_title) > 40: offset = Pt(60)
                        elif len(chart_title) > 54: offset = Pt(70)
                    else:
                        # 中文判断逻辑
                        if len(chart_title) > 20: offset = Pt(80)
                        elif len(chart_title) > 10: offset = Pt(30)
                    
                    
                    t_left = cfg['left_base'] - offset
                    
                    textbox = slide.shapes.add_textbox(t_left, cfg['top'], cfg['width'], cfg['height'])
                    run = textbox.text_frame.paragraphs[0].add_run()
                    run.text = chart_title
                    
                    # 3. 设置字体样式：如果是英文，强制使用 Arial Narrow
                    # 注意："Arial Narrow" 是字体名，bold=True 是加粗
                    font_to_use = "Arial Narrow" if is_english_mode else cfg['font_name']
                    
                    self._set_text_style(run, font_name=font_to_use, size=cfg['size'], bold=True)

            # B. 添加资料来源
            source_text = ""
            offset_s = Pt(0)
            
            if len(parts) >= 2:
                # 假设最后一部分是来源
                raw_source = parts[-1].replace('，', ' ')
                source_text = re.sub(r'\s+', ' ', raw_source).strip()
                if len(source_text) > 11 and self.language == "en": offset_s = Pt(70)
                elif len(source_text) > 13 and self.language != "cn": offset_s = Pt(40)
                else: offset_s = Pt(20)
            else:
                source_text = parts[0]
                offset_s = -Pt(15)
            
            source_text = source_text.strip('_ ')
            
            if source_text:
                cfg_s = config.ANNOTATION_CONFIG['source']
                s_left = cfg_s['left_base'] - offset_s
                
                textbox = slide.shapes.add_textbox(s_left, cfg_s['top'], cfg_s['width'], cfg_s['height'])
                # 确保每次都是新段落
                p = textbox.text_frame.paragraphs[0]
                if not p.runs: p.add_run()
                if self.language == "en":
                    source_text = GoogleTranslator(source='auto', target='en').translate(source_text)
               
                    print(f"Adding source annotation in English: {source_text}")
                    p.runs[0].text = f"Source: {source_text}"
                else:
                    source_text =self.translate_with_glossary(source_text)
                    p.runs[0].text = f"资料来源：{source_text}"
                
                self._set_text_style(p.runs[0], font_name=cfg_s['font_name'], size=cfg_s['size'], color=config.COLOR_GRAY)

        except Exception as e:
            logging.warning(f"添加图片注释失败: {e}")

    # --- 页面生成方法 ---

    def create_cover(self):
        """生成封面"""
        slide = self.prs.slides[0]
        doc = self.data.get("document", {})
        
        try:
            # 假设占位符顺序：0=日期, 1=作者, 2=标题
            slide.shapes[0].text = datetime.now().strftime("%Y-%m-%d")
            
           # 1. 设置作者 (Shape 1)
            # -------------------------------------------------------
            # 第一步：先获取文本框
            tf_author = slide.shapes[1].text_frame
            
            # 第二步：先填充文字内容 (这步会重置样式，所以必须先做)
            tf_author.paragraphs[0].text = doc.get("author", "")

            # 第三步：获取 Run 并设置样式 (加粗 + 微软雅黑)
            if tf_author.paragraphs[0].runs:
                self._set_text_style(
                    tf_author.paragraphs[0].runs[0], 
                    size=35, 
                    bold=True,                      # 加粗
                    color=config.COLOR_LIGHT_BLUE, 
                    font_name='Microsoft YaHei'      # <--- 这里指定字体
                )

            # -------------------------------------------------------
            # 2. 设置标题 (Shape 2)
            # -------------------------------------------------------
            # 第一步：先获取文本框
            tf_title = slide.shapes[2].text_frame
            
            # 第二步：先填充文字内容
            tf_title.paragraphs[0].text = doc.get("title", "")

            # 第三步：获取 Run 并设置样式 (加粗 + 微软雅黑)
            if tf_title.paragraphs[0].runs:
                self._set_text_style(
                    tf_title.paragraphs[0].runs[0], 
                    size=44, 
                    bold=True,                       
                    color=config.COLOR_DARK_BLUE, 
                    font_name='Microsoft YaHei'      # <--- 这里指定字体
                )
        except IndexError:
            logging.warning("封面页占位符索引不匹配")

    def create_summary(self):
        """生成摘要页"""
        slide = self.prs.slides[1]
        summary = self.data.get("executive_summary", {})
        if self.language == "en":
            size = 24
        else:
            size = 29.1
        # 标题
        if slide.shapes[0].has_text_frame:
            self._set_text_style(slide.shapes[0].text_frame.paragraphs[0].runs[0], 
                                 size=size, bold=True, color=config.COLOR_DARK_BLUE)

        # 表格
        try:
            table = slide.shapes[1].table
            cols = summary.get("columns", [])
            rows = summary.get("rows", [])

            for r_idx, row_data in enumerate(rows):
                if r_idx + 1 >= len(table.rows): break
                
                for c_idx, col_key in enumerate(cols):
                    # 计算表格实际列索引 (从索引1开始，即第二列)
                    table_col_idx = c_idx + 1
                    
                    if table_col_idx >= len(table.columns): break
                    
                    cell = table.cell(r_idx + 1, table_col_idx)
                    text = row_data.get(col_key, "")

                    # 1. 获取 TextFrame 并清空 (关键步骤)
                    tf = cell.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]

                    # 2. 添加 Run 并赋值
                    run = p.add_run()
                    run.text = text
                    
                    # 3. 计算样式逻辑
                    # 逻辑A: 投资逻辑字数多则变小
                    f_size = 10
                    
                    # 逻辑B: 判断是否为第二列 (index=1)，如果是则加粗
                    is_second_column = (table_col_idx == 1)

                    # 4. 应用样式
                    self._set_text_style(
                        run, 
                        size=f_size,
                        bold=is_second_column,       # <--- 这里控制加粗
                        font_name='Microsoft YaHei'  # 统一字体
                    )
                    
                    # 5. 设置居中 (仅第二列)
                    if is_second_column:
                        p.alignment = PP_ALIGN.CENTER
        
           

        except Exception as e:
            logging.error(f"摘要表格填充失败: {e}")

    def create_content_pages(self):
        """生成内容页"""
        # 清理旧幻灯片
        while len(self.prs.slides) > 2:
            rId = self.prs.slides._sldIdLst[2].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[2]

        layout = self.prs.slide_layouts[config.LAYOUT_IDX['content']]

        for i, content in  enumerate(self.data["content_slides"]):
            # 前6页尝试复用（如果有的话），否则新建
            if i < 6 and (2 + i) < len(self.prs.slides):
                slide = self.prs.slides[2 + i]
            else:
                slide = self.prs.slides.add_slide(layout)
            # 占位符
            if len(slide.placeholders) >= 2:
                title_ph = slide.placeholders[0]
                body_ph = slide.placeholders[10]
            else:
                # 如果占位符数量不足，尝试其他方法
                title_ph = None
                body_ph = None

            if title_ph and title_ph.has_text_frame:
                title_ph.text = content["title"]
                text_frame = title_ph.text_frame
                if text_frame.paragraphs and len(text_frame.paragraphs) > 0:
                    paragraph = text_frame.paragraphs[0]
                    if paragraph.runs and len(paragraph.runs) > 0:
                        run = paragraph.runs[0]
                        run.font.name = '华文细黑'
                        run.font.size = Pt(24.1)
                        run.font.bold = True
                        run.font.color.rgb = config.COLOR_DARK_BLUE

            # 正文
            if body_ph and body_ph.has_text_frame:
                # 添加项目符号
                bullets = content["bullets"]
                body_text = ""
                for bullet in bullets:
                    body_text += f"{bullet}\n"
                body_ph.text = body_text
                # 根据字数调整字体大小
                total_chars = len(body_text)
                font_size = Pt(14)
                
                for p in body_ph.text_frame.paragraphs:
                    for run in p.runs:
                        if self.language == "en":

                            run.font.size = Pt(12)
                        else:
                            run.font.size = Pt(14)
                print(f"Slide {i+1} Body Text added")

            # 图片
            matched_image_path = self._find_matching_image(content["title"], images_dir=self.images_dir)
            if matched_image_path:
                self._fill_image(slide, matched_image_path)
                # 内容页总是添加注释 
                self._add_image_annotations(slide, matched_image_path)

            self._remove_all_picture_placeholders(slide)
            print("-" * 50)

    def _format_title_placeholder(self, placeholder, text):
        """设置标题占位符的字体格式"""
        if placeholder and placeholder.has_text_frame:
            placeholder.text = text
            text_frame = placeholder.text_frame
            
            # 清除原有格式（如果有多个段落）
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = '华文细黑'
                    run.font.size = Pt(24.1)
                    run.font.bold = True
                    run.font.color.rgb = config.COLOR_DARK_BLUE
            
            # 确保至少有一个段落
            if text_frame.paragraphs and len(text_frame.paragraphs) > 0:
                paragraph = text_frame.paragraphs[0]
                # 确保至少有一个run
                if not paragraph.runs:
                    paragraph.text = text
                    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                else:
                    run = paragraph.runs[0]
                
                # 应用字体格式
                run.font.name = '华文细黑'
                run.font.size = Pt(24.1)
                run.font.bold = True
                run.font.color.rgb = config.COLOR_DARK_BLUE
    def create_image_slide(self, topic):
        """生成纯图页"""
        layout = self.prs.slide_layouts[config.LAYOUT_IDX['image_only']]
        slide = self.prs.slides.add_slide(layout)
        
        # 设置标题并应用字体格式
        title_ph = slide.placeholders[0]
        if title_ph and title_ph.has_text_frame:
            title_ph.text = topic
            text_frame = title_ph.text_frame
            
            # 清除原有格式（如果有多个段落）
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = '华文细黑'
                    run.font.size = Pt(24.1)
                    run.font.bold = True
                    run.font.color.rgb = config.COLOR_DARK_BLUE
            
            # 确保至少有一个段落
            if text_frame.paragraphs and len(text_frame.paragraphs) > 0:
                paragraph = text_frame.paragraphs[0]
                # 确保至少有一个run
                if not paragraph.runs:
                    paragraph.text = topic
                    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                else:
                    run = paragraph.runs[0]
                
                # 应用字体格式
                run.font.name = '华文细黑'
                run.font.size = Pt(24.1)
                run.font.bold = True
                run.font.color.rgb = config.COLOR_DARK_BLUE
        
        # 查找并添加图片
        img_path = self._find_matching_image(topic, images_dir=self.images_dir)
        if img_path:
            self._fill_image(slide, img_path)
            
            # 仅资金流加注释
            if topic == "资金流" or topic == "Fund Flow":
                self._add_image_annotations(slide, img_path)
        
        self._remove_all_picture_placeholders(slide)

    def create_contact_page(self):
        """生成联系页"""
        layout = self.prs.slide_layouts[config.LAYOUT_IDX['contact']]
        slide = self.prs.slides.add_slide(layout)
        # 根据语言设置标题文本
        title_text = "Contact Us" if self.language == "en" else "联系我们"
        
        # 设置标题并应用字体格式（使用公共方法）
        title_ph = slide.placeholders[0]
        self._format_title_placeholder(title_ph, title_text)
        # 主文本框 (索引15)
        try:
            if self.location == "香港/Hong Kong":
                # 1. 获取配置
                cfg = config.ANNOTATION_CONFIG['contact_info']
                
                # 2. 创建文本框 (位置：Left=1.06cm, Top=2.96cm)
                textbox = slide.shapes.add_textbox(cfg['left'], cfg['top'], cfg['width'], cfg['height'])
                tf = textbox.text_frame
                tf.word_wrap = True  # 允许自动换行

                # --- 第一行：网址 ---
                p1 = tf.paragraphs[0]
                run1 = p1.add_run()
                run1.text = "www.fargowealth.com"
                # 应用样式 (加粗，蓝色或黑色根据 config)
                self._set_text_style(run1, font_name=cfg['font_name'], size=cfg['size'], bold=False)
                run1.font.underline = True # 网址加下划线

                # --- 第二行：声明文字 ---
                p2 = tf.add_paragraph()
                if self.language == "en":
                    p2.text = "Asset management services are provided by Gentlemen Capital Limited, a subsidiary of Fargo Wealth Group."
                else:
                    p2.text = "资产管理服务由集团旗下公司绅士资本提供"
                
                # 为第二行应用样式 (通常声明文字不加粗，字号可以稍微调小一点)
                if p2.runs:
                    self._set_text_style(p2.runs[0], font_name=cfg['font_name'], size=cfg['size'], bold=False)

            else:
                # 非香港地区，移除所有文本占位符
                self._remove_all_text_placeholders(slide)

        except KeyError:
            pass
        if self.language == "en":
            print("Adding contact info in English")
            CONTACT_ADDRESSES = config.CONTACT_ADDRESSES_en
        else:
            CONTACT_ADDRESSES = config.CONTACT_ADDRESSES
        # 地址列表
        for idx, text in CONTACT_ADDRESSES.items():
            try:
                tf = slide.placeholders[idx].text_frame
                tf.clear()
                lines = text.split("\n")
                for i, line in enumerate(lines):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = line
                    for run in p.runs:
                        # 城市名加粗逻辑
                        if i == 0: run.font.bold = True
                        run.font.size = Pt(8)
            except KeyError:
                continue

    def create_disclaimer_pages(self):
        """生成免责页"""
        texts = config.DISCLAIMER_TEXTS.get(self.location, config.DISCLAIMER_TEXTS["default"])

        # 中文页
        if texts.get("cn"):
            layout = self.prs.slide_layouts[config.LAYOUT_IDX['disclaimer_cn']]
            slide = self.prs.slides.add_slide(layout)
            slide.placeholders[0].text = "免责声明"
            self._set_text_style(slide.placeholders[0].text_frame.paragraphs[0].runs[0], 
                                 size=29.1, bold=True, color=config.COLOR_DARK_BLUE)
            
            body = slide.placeholders[12]
            body.text = "\n".join(texts["cn"])
            for p in body.text_frame.paragraphs:
                for run in p.runs:
                    self._set_text_style(run, size=12)

        # 英文页 (非大陆)
        if self.location != "中国大陆" and texts.get("en"):
            layout = self.prs.slide_layouts[config.LAYOUT_IDX['disclaimer_en']]
            slide = self.prs.slides.add_slide(layout)
            slide.placeholders[0].text = "Disclaimer"
            self._set_text_style(slide.placeholders[0].text_frame.paragraphs[0].runs[0], 
                                 size=29.1, bold=True, color=config.COLOR_DARK_BLUE)
            
            body = slide.placeholders[12]
            body.text = "\n".join(texts["en"])
            for p in body.text_frame.paragraphs:
                for run in p.runs:
                    self._set_text_style(run, font_name='Arial', size=9)

    def run(self,output_path):
        """
        执行全流程
        :param output_path: 完整的输出文件路径 (包含目录和文件名)
        :return: Boolean (True 表示成功, False 表示失败)
        """
        logging.info(f"开始生成 PPT - 地点: {self.location}")
        # logging.info(f"开始生成 PPT - 地点: {self.location}")
        
        # --- 调试代码：查看当前 self.language 到底是什么 ---
        logging.info(f"DEBUG: 当前对象存储的语言为: '{self.language}'")
        
        try:
            # 1. 加载数据
            self.load_resources()
         
            # 2. 按顺序创建页面
            self.create_cover()              # 封面
            self.create_summary()            # 摘要
            self.create_content_pages()      # 核心内容
            
            # 图片页 (根据 key 查找图片)
            if self.language == "en":
                logging.info("添加英文版图片页")
                self.create_image_slide("Top Picks - Equities")
                self.create_image_slide("Top Picks - Bonds")
                self.create_image_slide("Fund Flow")
            else:
                self.create_image_slide("个股精选")
                self.create_image_slide("个债精选")
                self.create_image_slide("资金流")

            self.create_contact_page()       # 封底/联系
            self.create_disclaimer_pages()   # 免责声明
            
            # 3. 确保目标目录存在 (虽然 main.py 做了，这里再做一次保险)
            output_dir = os.path.dirname(output_path)
            if output_dir:
                os.makedirs(output_dir, exist_ok=True)
            
            # 4. 删除旧文件 (防止权限报错)
            if os.path.exists(output_path):
                try:
                    os.remove(output_path)
                except PermissionError:
                    logging.error(f"文件被占用，无法覆盖: {output_path}")
                    return False

            # 5. 保存文件
            self.prs.save(output_path)
            logging.info(f"PPT 生成完成: {output_path} (共 {len(self.prs.slides)} 页)")
            
            # ----------------------------------------
            # 【关键】必须返回 True，否则 main.py 会认为失败
            # ----------------------------------------
            return True

        except Exception as e:
            logging.error(f"PPT 生成过程中发生异常: {str(e)}")
            import traceback
            logging.error(traceback.format_exc())
            return False


# ================= 对外接口函数 =================

def generate_ppt_from_json(json_path, template_path, output_filename, location_name, images_dir, language="cn"):
    try:
        # 1. Must load the JSON first
        with open(json_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
            
        # 2. Pass data (dict), not the path (string)
        generator = PPTGenerator(data, template_path, images_dir, location_name, language)
        return generator.run(output_filename)
    except Exception as e:
        logging.error(f"PPT Generation Error: {e}")
        return False
# 示例调用
if __name__ == "__main__":
    # 配置
    IMAGES_DIR = "input_articles/20260116/images_20260116"
    DATA_JSON = "final_investment_report.json"
    
    # 交互输入
    print("请选择需要PPT的地点：\n1. 中国大陆\n2. 香港\n3. 新加坡")
    choice = input("请输入数字 (1/2/3): ").strip()
    
    LOCATION_MAP = {
        "1": ("中国大陆", "template/AI PPT v3.pptx"),
        "2": ("香港", "template/AI PPT v2.pptx"),
        "3": ("新加坡", "template/AI PPT v2.pptx")
    }
    
    if choice not in LOCATION_MAP:
        print("输入无效，程序退出")
        exit(1)
        
    loc_name, template_file = LOCATION_MAP[choice]
    output_file = f"AI_PPT_generated_{loc_name}.pptx"
    
    generate_ppt_from_json(DATA_JSON, template_file, output_file, loc_name, IMAGES_DIR)
